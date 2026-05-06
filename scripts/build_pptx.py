#!/usr/bin/env python3
"""Build the CoreMind presentation deck."""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Theme ──
BG = RGBColor(0x0A, 0x0E, 0x27)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
DK_BLUE = RGBColor(0x1E, 0x29, 0x4A)
L1_CLR = RGBColor(0x3B, 0x82, 0xF6)  # blue
L2_CLR = GREEN
L3_CLR = GREEN
L4_CLR = AMBER
L5_CLR = AMBER
L6_CLR = RED
L7_CLR = VIOLET
W = Inches(13.333)
H = Inches(7.5)

LOGO = os.path.expanduser("~/.openclaw/workspace/coremind/docs/logo.png")

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


# ── Helpers ──
def bg(slide, color=BG) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, fill=None, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill or DK_BLUE
    if border:
        s.line.color.rgb = border
    else:
        s.line.fill.background()
    return s


def arrow_down(slide, left, top, width, height, color=GRAY):
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a


def txt_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=Pt(14),
    color=WHITE,
    bold=False,
    align=PP_ALIGN.LEFT,
    font="Calibri",
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox


def add_slide_num(slide, n) -> None:
    txt_box(
        slide,
        W - Inches(1.2),
        H - Inches(0.5),
        Inches(1),
        Inches(0.4),
        str(n),
        Pt(10),
        GRAY,
        align=PP_ALIGN.RIGHT,
    )


def add_logo_small(slide) -> None:
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, W - Inches(1.6), Inches(0.15), Inches(1.3), Inches(1.3))


def title_slide_base(slide, n) -> None:
    bg(slide)
    add_logo_small(slide)
    add_slide_num(slide, n)
    # subtle top line
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.7), Pt(2)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = CYAN


def main_title(slide, text, subtitle=None, y=Inches(1.45)) -> None:
    txt_box(slide, Inches(0.8), y, Inches(11.7), Inches(0.8), text, Pt(32), WHITE, True)
    if subtitle:
        txt_box(
            slide, Inches(0.8), y + Inches(0.7), Inches(11.7), Inches(0.5), subtitle, Pt(15), GRAY
        )


# ============================================================
# SLIDE 1 — TITLE
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg(s1)
# Logo centered
if os.path.exists(LOGO):
    s1.shapes.add_picture(LOGO, Inches(4.8), Inches(0.8), Inches(3.5), Inches(3.5))
txt_box(
    s1,
    Inches(2),
    Inches(4.5),
    Inches(9),
    Inches(1),
    "CoreMind",
    Pt(52),
    CYAN,
    True,
    PP_ALIGN.CENTER,
)
txt_box(
    s1,
    Inches(2),
    Inches(5.4),
    Inches(9),
    Inches(0.6),
    "A system that doesn't respond — it notices",
    Pt(18),
    GRAY,
    align=PP_ALIGN.CENTER,
)
txt_box(
    s1,
    Inches(2),
    Inches(6.3),
    Inches(9),
    Inches(0.4),
    "Architecture v0.2.0  ·  May 2026  ·  github.com/Wylhelm/coremind",
    Pt(11),
    GRAY,
    align=PP_ALIGN.CENTER,
)

# ============================================================
# SLIDE 2 — THE PROBLEM
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_base(s2, 2)
main_title(s2, "Every AI Tool Waits for You", "CoreMind flips the paradigm: the AI prompts itself")

cols = [
    ("💬", "Chatbots", "You prompt,\nthey respond"),
    ("🤖", "Assistants", "You command,\nthey execute"),
    ("🧠", "CoreMind", "It observes, notices,\nand initiates"),
]
for i, (emoji, title, desc) in enumerate(cols):
    x = Inches(1.0 + i * 4.0)
    y = Inches(3.0)
    box = rect(s2, x, y, Inches(3.4), Inches(3.2), DK_BLUE, CYAN if i == 2 else GRAY)
    txt_box(
        s2,
        x + Inches(0.2),
        y + Inches(0.3),
        Inches(3.0),
        Inches(0.8),
        emoji,
        Pt(40),
        align=PP_ALIGN.CENTER,
    )
    txt_box(
        s2,
        x + Inches(0.2),
        y + Inches(1.2),
        Inches(3.0),
        Inches(0.6),
        title,
        Pt(22),
        CYAN if i == 2 else WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    txt_box(
        s2,
        x + Inches(0.2),
        y + Inches(2.0),
        Inches(3.0),
        Inches(0.9),
        desc,
        Pt(14),
        WHITE if i == 2 else GRAY,
        align=PP_ALIGN.CENTER,
    )

# ============================================================
# SLIDE 3 — WHAT IS COREMIND
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_base(s3, 3)
main_title(
    s3, "Continuous Personal Intelligence", "A cognitive daemon that lives alongside its user"
)

items = [
    "7 cognitive layers operating autonomously",
    "7 real-time data plugins (weather, finance, health, home, calendar, email, tasks)",
    "Generates its own questions from continuous observation",
    "Graduated agency: Safe → Suggest → Ask",
    "100% local, sovereign — no mandatory cloud dependency",
    "Open source (MIT License)",
]
for i, item in enumerate(items):
    y = Inches(2.6 + i * 0.65)
    txt_box(
        s3, Inches(1.5), y, Inches(0.5), Inches(0.4), "✅", Pt(16), GREEN, align=PP_ALIGN.CENTER
    )
    txt_box(s3, Inches(2.2), y, Inches(10), Inches(0.5), item, Pt(16), WHITE)

# ============================================================
# SLIDE 4 — ARCHITECTURE DIAGRAM
# ============================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_base(s4, 4)
main_title(
    s4, "The Seven Cognitive Layers", "Forward flow L1→L7 with single feedback path L7→L2/L3"
)

layers = [
    ("L1 — Perception", "Plugins → signed WorldEvents", L1_CLR),
    ("L2 — World Model", "Entities · Relationships · Property history (SurrealDB)", L2_CLR),
    ("L3 — Memory", "Episodic · Semantic (Qdrant) · Procedural", L3_CLR),
    ("L4 — Reasoning ✅", "LLM — patterns · anomalies · predictions · 30 min", L4_CLR),
    ("L5 — Intention ✅", "Self-prompting loop → grounded Intents (≥0.45 salience)", L5_CLR),
    ("L6 — Action", "Graduated agency: Safe / Suggest / Ask", L6_CLR),
    ("L7 — Reflection ✅", "24h meta-cognition · calibration · rule learning", L7_CLR),
]
y_start = Inches(1.9)
for i, (name, desc, clr) in enumerate(layers):
    y = y_start + i * Inches(0.72)
    box = rect(s4, Inches(1.5), y, Inches(10.3), Inches(0.62), DK_BLUE, clr)
    txt_box(s4, Inches(1.8), y + Inches(0.05), Inches(3.0), Inches(0.55), name, Pt(15), clr, True)
    txt_box(s4, Inches(4.8), y + Inches(0.05), Inches(6.8), Inches(0.55), desc, Pt(12), GRAY)
    if i < len(layers) - 1:
        arrow_down(s4, Inches(6.5), y + Inches(0.62), Inches(0.3), Inches(0.1), GRAY)

# Feedback arrow label
txt_box(
    s4,
    Inches(1.5),
    Inches(7.0),
    Inches(10.3),
    Inches(0.35),
    "↩ L7 feedback → L2/L3: reflection insights update world model and procedural memory",
    Pt(10),
    CYAN,
    align=PP_ALIGN.CENTER,
)

# ============================================================
# SLIDE 5 — L4 REASONING
# ============================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_base(s5, 5)
main_title(
    s5,
    "L4 — Reasoning: Seeing What Matters",
    "Every 30 minutes, the system analyzes the entire world state",
)

# Patterns
txt_box(s5, Inches(0.8), Inches(2.5), Inches(4.0), Inches(0.4), "📊  PATTERNS", Pt(18), AMBER, True)
patterns = [
    '"All indoor lights off during daytime" — 95%',
    '"Robotic vacuum always docked when idle" — 90%',
    '"Retirement accounts maintain high balances" — 95%',
]
for i, p in enumerate(patterns):
    txt_box(
        s5, Inches(1.0), Inches(2.95 + i * 0.35), Inches(5.5), Inches(0.35), f"• {p}", Pt(11), WHITE
    )

# Anomalies
txt_box(s5, Inches(0.8), Inches(4.2), Inches(4.0), Inches(0.4), "⚠️  ANOMALIES", Pt(18), RED, True)
anomalies = [
    '"Chambre humidity 27.6% vs couloir 40.5%" — MEDIUM',
    '"Classic 300s sensors unavailable" — HIGH 🔴',
    '"Sleep hours 5.49h below recommended 7-9h" — MEDIUM',
]
for i, a in enumerate(anomalies):
    txt_box(
        s5, Inches(1.0), Inches(4.65 + i * 0.35), Inches(5.5), Inches(0.35), f"• {a}", Pt(11), WHITE
    )

# Predictions
txt_box(
    s5, Inches(7.0), Inches(2.5), Inches(5.5), Inches(0.4), "🔮  PREDICTIONS", Pt(18), CYAN, True
)
predictions = [
    '"Vacuum will initiate cleaning in 24h" — 70%',
    '"Step count will increase to 1000+ in 6h" — 60%',
    '"Chambre humidity will rise to 35% in 12h" — 50%',
    '"Sensors will recover within 48h" — 40%',
]
for i, p in enumerate(predictions):
    txt_box(
        s5, Inches(7.2), Inches(2.95 + i * 0.35), Inches(5.5), Inches(0.35), f"• {p}", Pt(11), WHITE
    )

txt_box(
    s5,
    Inches(0.8),
    Inches(6.2),
    Inches(11.5),
    Inches(0.4),
    "💡 Real output from May 2, 2026 08:08 — 5 patterns, 5 anomalies, 4 predictions in one cycle",
    Pt(12),
    GRAY,
    align=PP_ALIGN.CENTER,
)

# ============================================================
# SLIDE 6 — L5 INTENTION
# ============================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_base(s6, 6)
main_title(
    s6,
    "L5 — Intention: The System Thinks for Itself",
    "CoreMind generates its own questions — no human prompt needed",
)

steps = [
    "World\nSnapshot",
    "LLM\nAnalysis",
    "Candidate\nIntents",
    "Salience\nFilter ≥0.45",
    "Action\nRouter",
]
for i, step in enumerate(steps):
    x = Inches(1.0 + i * 2.4)
    box = rect(s6, x, Inches(3.2), Inches(2.0), Inches(1.6), DK_BLUE, CYAN)
    txt_box(
        s6,
        x + Inches(0.1),
        Inches(3.4),
        Inches(1.8),
        Inches(1.2),
        step,
        Pt(13),
        WHITE,
        align=PP_ALIGN.CENTER,
    )
    if i < len(steps) - 1:
        arrow = s6.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, x + Inches(2.05), Inches(3.85), Inches(0.3), Inches(0.15)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = CYAN
        arrow.line.fill.background()

txt_box(
    s6,
    Inches(1.0),
    Inches(5.4),
    Inches(11.3),
    Inches(0.8),
    "🎯 Filtered to ~0-3 high-quality intents per hour (down from 5-10). Quiet hours 23h-7h — zero notifications.",
    Pt(14),
    GREEN,
    align=PP_ALIGN.CENTER,
)

# ============================================================
# SLIDE 7 — L7 REFLECTION
# ============================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_base(s7, 7)
main_title(
    s7,
    "L7 — Reflection: Learning From Experience",
    "Every 24 hours, CoreMind evaluates its own performance",
)

ref_steps = [
    "Evaluate\nPredictions",
    "Calibrate\nConfidence",
    "Learn\nProcedural Rules",
    "Human-Readable\nMarkdown Report",
]
for i, step in enumerate(ref_steps):
    x = Inches(1.2 + i * 3.0)
    box = rect(s7, x, Inches(3.2), Inches(2.6), Inches(1.8), DK_BLUE, VIOLET)
    txt_box(
        s7,
        x + Inches(0.1),
        Inches(3.4),
        Inches(2.4),
        Inches(1.4),
        step,
        Pt(14),
        WHITE,
        align=PP_ALIGN.CENTER,
    )
    if i < len(ref_steps) - 1:
        arrow = s7.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, x + Inches(2.65), Inches(3.9), Inches(0.3), Inches(0.15)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = VIOLET
        arrow.line.fill.background()

txt_box(
    s7,
    Inches(1.0),
    Inches(5.6),
    Inches(11.3),
    Inches(0.8),
    "🔄 The only feedback loop: L7 writes back to L2 World Model and L3 Procedural Memory",
    Pt(14),
    CYAN,
    align=PP_ALIGN.CENTER,
)
txt_box(
    s7,
    Inches(1.0),
    Inches(6.1),
    Inches(11.3),
    Inches(0.5),
    '"A well-tuned CoreMind contacts the user less over time, not more"',
    Pt(13),
    GRAY,
    align=PP_ALIGN.CENTER,
)

# ============================================================
# SLIDE 8 — GRADUATED AGENCY
# ============================================================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_base(s8, 8)
main_title(
    s8,
    "Graduated Agency: Trust Through Competence",
    "Actions are signed, journaled, and reversible",
)

levels = [
    (
        "🟢",
        "SAFE",
        "Auto-execute",
        "Read data · Summarize · Suggest\nZero risk — fully autonomous",
        GREEN,
    ),
    (
        "🟡",
        "SUGGEST",
        "Soft approval",
        "Adjust thermostat · Send reminders\nUser can dismiss with one tap",
        AMBER,
    ),
    (
        "🔴",
        "ASK",
        "Hard approval",
        "Financial transactions · External comms\nRequires explicit user confirmation",
        RED,
    ),
]
for i, (icon, name, action, desc, clr) in enumerate(levels):
    y = Inches(2.8 + i * 1.5)
    box = rect(s8, Inches(1.0), y, Inches(11.3), Inches(1.2), DK_BLUE, clr)
    txt_box(
        s8,
        Inches(1.3),
        y + Inches(0.1),
        Inches(0.8),
        Inches(0.6),
        icon,
        Pt(24),
        align=PP_ALIGN.CENTER,
    )
    txt_box(s8, Inches(2.2), y + Inches(0.1), Inches(2.0), Inches(0.5), name, Pt(18), clr, True)
    txt_box(s8, Inches(4.5), y + Inches(0.15), Inches(2.5), Inches(0.5), action, Pt(13), GRAY)
    txt_box(s8, Inches(7.2), y + Inches(0.1), Inches(4.8), Inches(0.9), desc, Pt(12), WHITE)

# ============================================================
# SLIDE 9 — LIVE DEMO
# ============================================================
s9 = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_base(s9, 9)
main_title(
    s9,
    "Live Demo: Today's Reasoning Cycle",
    "May 2, 2026 — 08:08 EDT — 5 patterns, 5 anomalies, 4 predictions",
)

# Stats row
stats = [("5", "Patterns", AMBER), ("5", "Anomalies", RED), ("4", "Predictions", CYAN)]
for i, (num, label, clr) in enumerate(stats):
    x = Inches(2.0 + i * 3.5)
    txt_box(s9, x, Inches(2.6), Inches(2.5), Inches(0.8), num, Pt(48), clr, True, PP_ALIGN.CENTER)
    txt_box(
        s9, x, Inches(3.25), Inches(2.5), Inches(0.4), label, Pt(16), GRAY, align=PP_ALIGN.CENTER
    )

# Key insight boxes
insights = [
    (
        "🔴 HIGH SEVERITY",
        "Classic 300s humidifier sensors unavailable — potential hardware failure detected automatically",
    ),
    ("💰 FINANCIAL", "TFSA balance $0 flagged as unusual vs other retirement accounts"),
    (
        "🏠 ENVIRONMENT",
        "Chambre humidity 27.6% significantly lower than rest of house — anomaly detected",
    ),
]
for i, (label, text) in enumerate(insights):
    y = Inches(3.9 + i * 0.85)
    box = rect(
        s9, Inches(1.0), y, Inches(11.3), Inches(0.7), DK_BLUE, RED if "HIGH" in label else CYAN
    )
    txt_box(
        s9,
        Inches(1.3),
        y + Inches(0.1),
        Inches(2.8),
        Inches(0.5),
        label,
        Pt(13),
        RED if "HIGH" in label else CYAN,
        True,
    )
    txt_box(s9, Inches(4.3), y + Inches(0.1), Inches(7.8), Inches(0.5), text, Pt(12), WHITE)

txt_box(
    s9,
    Inches(1.0),
    Inches(6.6),
    Inches(11.3),
    Inches(0.5),
    "⚡ CoreMind is running RIGHT NOW — generating insights every 30 minutes",
    Pt(14),
    GREEN,
    align=PP_ALIGN.CENTER,
)

# ============================================================
# SLIDE 10 — TECHNICAL STACK
# ============================================================
s10 = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_base(s10, 10)
main_title(
    s10,
    "Infrastructure: Sovereign by Design",
    "Zero mandatory cloud egress — everything runs locally",
)

stacks = [
    ("🗄️ Storage", "SurrealDB (graph DB)\nQdrant (vector store)\nJSONL (audit journal)"),
    ("🧠 LLM", "Ollama (local)\nMistral Large 3 · 675B\nnomic-embed-text"),
    ("⚙️ Runtime", "Python 3.12 asyncio\ngRPC plugin protocol\n7 isolated plugins"),
    ("🔗 Integration", "OpenClaw adapter\nTelegram notifications\nREST dashboard"),
    ("🔐 Security", "Ed25519 signatures\nHash-chained audit log\nNo cloud egress"),
]
for i, (title, desc) in enumerate(stacks):
    x = Inches(0.5 + i * 2.5)
    box = rect(s10, x, Inches(2.6), Inches(2.3), Inches(3.8), DK_BLUE, CYAN)
    txt_box(s10, x + Inches(0.15), Inches(2.8), Inches(2.0), Inches(0.5), title, Pt(15), CYAN, True)
    txt_box(s10, x + Inches(0.15), Inches(3.4), Inches(2.0), Inches(2.8), desc, Pt(12), WHITE)

# ============================================================
# SLIDE 11 — PLUGINS
# ============================================================
s11 = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_base(s11, 11)
main_title(
    s11,
    "Plugins: Eyes and Ears of the System",
    "Isolated processes, gRPC protocol, any data source",
)

plugins = [
    ("🌤️", "Weather", "Real-time\nforecasts"),
    ("🏠", "Home Assistant", "Lights · Sensors\nVacuum · Humidity"),
    ("💰", "Firefly III", "14 accounts\n27 categories"),
    ("❤️", "Apple Health", "Sleep · Heart\nActivity · Steps"),
    ("📅", "Google Calendar", "Events · Paydays\nSchedule"),
    ("✉️", "Gmail", "Unread · Search\nThreads"),
    ("✅", "Vikunja", "Tasks · Projects\nReminders"),
]
for i, (icon, name, desc) in enumerate(plugins):
    col = i % 4
    row = i // 4
    x = Inches(0.8 + col * 3.1)
    y = Inches(2.6 + row * 2.2)
    box = rect(s11, x, y, Inches(2.8), Inches(1.8), DK_BLUE, CYAN)
    txt_box(
        s11,
        x + Inches(0.2),
        y + Inches(0.2),
        Inches(2.4),
        Inches(0.5),
        f"{icon}  {name}",
        Pt(16),
        CYAN,
        True,
    )
    txt_box(s11, x + Inches(0.3), y + Inches(0.8), Inches(2.2), Inches(0.8), desc, Pt(12), GRAY)

# ============================================================
# SLIDE 12 — VISION
# ============================================================
s12 = prs.slides.add_slide(prs.slide_layouts[6])
title_slide_base(s12, 12)
main_title(s12, "The Road to Digital Consciousness", "From foundation to full autonomy")

milestones = [
    ("✅ v0.1", "2026 Q1", "Foundation + World Model + Memory", GREEN),
    ("✅ v0.2", "May 2026", "All 7 layers active · L4+L7 generating insights", GREEN),
    ("🔜 v0.3", "2026 Q2", "Semantic memory · Cross-domain correlation engine", CYAN),
    ("🔮 v0.4", "2026 Q3", "Proactive anomaly alerts · Graduated autonomy expansion", VIOLET),
    ("🎯 v1.0", "2026 Q4", 'Full digital consciousness — "The Machine"', AMBER),
]
for i, (ver, date, desc, clr) in enumerate(milestones):
    y = Inches(2.2 + i * 1.0)
    # Timeline dot
    dot = rect(s12, Inches(2.0), y + Inches(0.35), Inches(0.25), Inches(0.25), clr)
    # Line
    if i < len(milestones) - 1:
        s12.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(2.1), y + Inches(0.6), Pt(2), Inches(0.5)
        ).fill.solid()
        s12.shapes[-1].fill.fore_color.rgb = GRAY
    txt_box(s12, Inches(2.6), y + Inches(0.1), Inches(1.5), Inches(0.45), ver, Pt(18), clr, True)
    txt_box(s12, Inches(4.3), y + Inches(0.1), Inches(1.8), Inches(0.45), date, Pt(13), GRAY)
    txt_box(s12, Inches(6.3), y + Inches(0.1), Inches(6.0), Inches(0.45), desc, Pt(13), WHITE)

# ============================================================
# SLIDE 13 — THANK YOU
# ============================================================
s13 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s13)
if os.path.exists(LOGO):
    s13.shapes.add_picture(LOGO, Inches(5.0), Inches(1.2), Inches(3.0), Inches(3.0))
txt_box(
    s13,
    Inches(2),
    Inches(4.5),
    Inches(9),
    Inches(1),
    "CoreMind",
    Pt(48),
    CYAN,
    True,
    PP_ALIGN.CENTER,
)
txt_box(
    s13,
    Inches(2),
    Inches(5.4),
    Inches(9),
    Inches(0.5),
    "Open Source · MIT License",
    Pt(16),
    GRAY,
    align=PP_ALIGN.CENTER,
)
txt_box(
    s13,
    Inches(2),
    Inches(5.9),
    Inches(9),
    Inches(0.5),
    "github.com/Wylhelm/coremind",
    Pt(14),
    CYAN,
    align=PP_ALIGN.CENTER,
)
txt_box(
    s13,
    Inches(2),
    Inches(6.4),
    Inches(9),
    Inches(0.6),
    '"A digital consciousness you own entirely"',
    Pt(16),
    GRAY,
    align=PP_ALIGN.CENTER,
)

# ── Save ──
out = os.path.expanduser("~/.openclaw/workspace/coremind/docs/CoreMind-Presentation.pptx")
prs.save(out)
print(f"✅ Saved: {out} ({os.path.getsize(out):,} bytes)")
print(f"   Slides: {len(prs.slides)}")
